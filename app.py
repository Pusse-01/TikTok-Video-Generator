# streamlit_app.py
import io
import os
import re
import math
import zipfile
import tempfile
import shutil
from pathlib import Path
from typing import List, Tuple

import streamlit as st

# Video / text tools
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from scenedetect import open_video, SceneManager
from scenedetect.detectors import ContentDetector

from moviepy.editor import (
    VideoFileClip,
    CompositeVideoClip,
    ImageClip,
    concatenate_videoclips,
    vfx,
)

# -----------------------------
# --------- UI STYLE ----------
# -----------------------------
st.set_page_config(page_title="NP360 TikTok Generator", page_icon="🎬", layout="wide")

st.markdown("""
<style>
/* cleaner look */
.block-container {padding-top: 1.5rem; padding-bottom: 2rem; }
[data-testid="stSidebar"] {min-width: 360px; width: 360px;}
div.stProgress > div > div > div { background-image: none !important; }
.small-muted {color:#6b7280; font-size: 0.9rem;}
.badge {display:inline-block; padding:2px 8px; border-radius:999px; border:1px solid #e5e7eb; font-size:12px; margin-right:6px;}
hr {border-top:1px solid #e5e7eb;}
</style>
""", unsafe_allow_html=True)

EMOJI_RENDER_MODE = "skia"   # "skia" | "twemoji" | "pil"
EMOJI_FONT_PATH   = "assets/Fonts/NotoColorEmoji-Regular.ttf"  # put the font here

# -----------------------------
# ---------- HELPERS ----------
# -----------------------------

TIME_RE = re.compile(r"^\s*(\d+):(\d{1,2})(?:\.(\d{1,3}))?\s*$")  # mm:ss(.ms)

def parse_time_to_seconds(s: str) -> float:
    """Accepts 'mm:ss', 'mm:ss.ms', or raw seconds like '7.5'."""
    s = str(s).strip()
    m = TIME_RE.match(s)
    if m:
        mm = int(m.group(1))
        ss = int(m.group(2))
        ms = int(m.group(3)) if m.group(3) else 0
        return mm * 60 + ss + ms / (1000 if ms >= 100 else (10 if ms < 10 else 100))
    # fallback: plain float seconds
    return float(s)

def valid_hex_color(s: str) -> bool:
    return bool(re.fullmatch(r"#(?:[0-9a-fA-F]{6}|[0-9a-fA-F]{3})", s.strip()))

def clamp(v: float, lo: float, hi: float) -> float:
    return max(lo, min(hi, v))

def resolve_font_path(font_key: str, cfg: dict) -> str:
    """
    font_key can be 'default', 'cute', 'qq', or a path to a font file.
    Falls back to cfg['font_path'] if not found.
    """
    if not font_key:
        return cfg["font_path"]
    key = font_key.strip().lower()
    packs = cfg.get("font_packs", {})
    if key in packs and packs[key]:
        return packs[key]
    # treat as a direct path
    return font_key if Path(font_key).exists() else cfg["font_path"]

TAG_LINE_RE = re.compile(r"^\s*\[(?P<tags>[^\]]+)\]\s*(?P<text>.*)$")

def parse_tag_block(tags_raw: str) -> dict:
    """
    Supports:
      • time range: '0:03-0:06' or '3-6'
      • key=value pairs: pos=top|middle|bottom, font=cute|qq|<path>, color=#RRGGBB, size=72
    Returns a dict with any of: start, end, pos, font, color, size
    """
    info = {}
    parts = [p for p in tags_raw.strip().split() if p.strip()]
    for p in parts:
        # time range?
        if "-" in p and not ("=" in p):
            try:
                a, b = p.split("-", 1)
                info["start"] = parse_time_to_seconds(a)
                info["end"] = parse_time_to_seconds(b)
            except Exception:
                pass
            continue
        # key=value
        if "=" in p:
            k, v = p.split("=", 1)
            k = k.strip().lower()
            v = v.strip()
            if k == "pos" and v.lower() in {"top", "middle", "bottom"}:
                info["pos"] = v.lower()
            elif k == "font":  # 'default', 'cute', 'qq', or path
                info["font"] = v
            elif k == "color" and valid_hex_color(v):
                info["color"] = v
            elif k == "size":
                try:
                    info["size"] = int(v)
                except Exception:
                    pass
            elif k in {"start", "end"}:
                try:
                    info[k] = parse_time_to_seconds(v)
                except Exception:
                    pass
    return info

def parse_tagged_line(raw: str) -> tuple[dict, str]:
    """
    Returns (meta, text). meta may contain start/end/pos/font/color/size.
    If no tag block, meta is {} and text is the raw (trimmed).
    """
    m = TAG_LINE_RE.match(raw or "")
    if not m:
        return {}, (raw or "").strip()
    meta = parse_tag_block(m.group("tags"))
    text = (m.group("text") or "").strip()
    return meta, text

from dataclasses import dataclass

@dataclass
class SlideParseResult:
    cues: List[dict]   # each: {text, start?, dur?, pos?, font?, color?, size?}
    caption: str | None

def extract_struct_from_pptx(pptx_path: str, slide_keyword: str, default_pos: str = "middle") -> SlideParseResult:
    """
    Finds the first slide containing `slide_keyword` anywhere in its concatenated text.
    Parses each paragraph:
      - 'Copy:' → captured as TikTok caption (not on-screen)
      - CJK-heavy lines are skipped (headlines)
      - Tag format: [0:03-0:06 pos=top font=cute color=#FF61C0 size=72] Actual text
    Returns cues with optional explicit start/end; if both present we compute dur.
    """
    prs = Presentation(pptx_path)
    chosen = None
    for s in prs.slides:
        text = " ".join([sh.text for sh in s.shapes if hasattr(sh, "text")])
        if slide_keyword in text:
            chosen = s
            break

    if not chosen:
        st.warning(f"⚠️ No slide found containing keyword: '{slide_keyword}' for ")
        return SlideParseResult(cues=[], caption=None)

    cues: List[dict] = []
    caption = None

    for sh in chosen.shapes:
        # Skip title placeholders
        if getattr(sh, "placeholder_format", None) and sh.placeholder_format.type == 1:
            continue
        if hasattr(sh, "text_frame") and sh.text_frame:
            for p in sh.text_frame.paragraphs:
                t = (p.text or "").strip()
                if not t:
                    continue
                if t.lower().startswith("copy:"):
                    # caption line
                    cp = t[5:].strip()
                    caption = cp if cp else caption
                    continue
                if is_mostly_cjk(t):
                    # likely a large headline; ignore by default
                    continue

                meta, text = parse_tagged_line(t)
                if not text:
                    continue
                item = {
                    "text": text,
                    # defaults applied later
                }
                # attach meta
                item.update({k: v for k, v in meta.items() if v is not None})
                # default position if not supplied
                item["pos"] = (meta.get("pos") or default_pos or "middle").lower()
                # compute dur if start & end provided
                if "start" in item and "end" in item:
                    item["dur"] = max(0.05, float(item["end"]) - float(item["start"]))
                cues.append(item)

    # de-dup (by text + timing signature) while preserving order
    seen = set()
    uniq = []
    for c in cues:
        sig = (c["text"], round(c.get("start", -1), 3), round(c.get("dur", -1), 3), c.get("pos", ""))
        if sig in seen:
            continue
        seen.add(sig)
        uniq.append(c)

    return SlideParseResult(cues=uniq, caption=caption)


def is_mostly_cjk(t: str, frac=0.6) -> bool:
    """Checks whether the text is chinese/korean/japanese"""
    if not t:
        return False
    cjk = sum(1 for ch in t if '\u4e00' <= ch <= '\u9fff' or '\u3040' <= ch <= '\u30ff')
    return (cjk / max(1, len(t))) >= frac


def extract_lines_from_pptx(pptx_path: str, slide_keyword: str) -> List[str]:
    """Extract cleaned lines for a given slide keyword (e.g., the file name)."""
    prs = Presentation(pptx_path)
    chosen = None
    for s in prs.slides:
        text = " ".join([sh.text for sh in s.shapes if hasattr(sh, "text")])
        if slide_keyword in text:
            chosen = s
            break
    if not chosen:
        st.warning(f"⚠️ No slide found containing keyword: '{slide_keyword}'")
        return []

    lines = []
    for sh in chosen.shapes:
        # Skip slide title placeholders entirely
        if getattr(sh, "placeholder_format", None) and sh.placeholder_format.type == 1:
            continue
        if hasattr(sh, "text_frame") and sh.text_frame:
            for p in sh.text_frame.paragraphs:
                t = p.text.strip()
                if not t:
                    continue
                if t.lower().startswith("copy:"):
                    continue  # caption for TikTok, not on-screen
                if is_mostly_cjk(t):
                    continue  # ignore large Chinese headline
                lines.append(t)

    # De-dup keep order
    clean, seen = [], set()
    for l in lines:
        if l not in seen:
            seen.add(l)
            clean.append(l)
    return clean


def detect_scenes(video_path: str, threshold: float = 27.0) -> List[Tuple[float, float]]:
    """Return list of (start, end) seconds for scenes. Fallback to whole clip."""
    # with VideoFileClip(video_path) as c:
    #     return [(0.0, c.duration)]
    video = open_video(video_path)
    sm = SceneManager()
    sm.add_detector(ContentDetector(threshold=threshold))
    sm.detect_scenes(video)
    lst = sm.get_scene_list()
    if not lst:
        return [(0.0, video.duration.get_seconds())]
    return [(s.get_seconds(), e.get_seconds()) for (s, e) in lst]


def merge_scenes_for_readability(scenes, min_len=2.2, max_len=5.0):
    """Merge rapid cuts into readable beats; cap overly long beats."""
    if not scenes:
        return []
    beats = []
    cur_s, cur_e = scenes[0]
    for s, e in scenes[1:]:
        if (cur_e - cur_s) < min_len:
            cur_e = e
        else:
            beats.append((cur_s, cur_e))
            cur_s, cur_e = s, e
    beats.append((cur_s, cur_e))

    capped = []
    for s, e in beats:
        if e - s > max_len:
            capped.append((s, s + max_len))
        else:
            capped.append((s, e))
    return capped


def plan_cues_readable(lines, scenes, clip_duration, min_read=2.2, max_read=3.5, gap=0.12):
    """No overlaps, readable durations, skip micro-scenes intelligently."""
    beats = merge_scenes_for_readability(scenes, min_len=min_read, max_len=5.0) or [(0.0, clip_duration)]
    cues = []
    if not lines:
        return cues

    bi = 0
    for li, line in enumerate(lines):
        bstart, bend = beats[min(bi, len(beats)-1)]
        avail = max(0.0, bend - bstart - gap)
        dur = max(min_read, min(max_read, avail))
        start = bstart + gap

        # push right to avoid overlap
        if cues and start < (cues[-1]["start"] + cues[-1]["dur"] + gap):
            start = cues[-1]["start"] + cues[-1]["dur"] + gap
            if start + dur > bend:
                bi = min(bi + 1, len(beats)-1)
                bstart, bend = beats[bi]
                start = bstart + gap
                avail = max(0.0, bend - bstart - gap)
                dur = max(min_read, min(max_read, avail))

        cues.append({"text": line, "start": start, "dur": dur, "pos": "bottom"})
        bi = min(bi + 1, len(beats)-1)

    # clamp within clip
    for c in cues:
        if c["start"] + c["dur"] > clip_duration:
            c["dur"] = max(0.8, clip_duration - c["start"] - gap)
    return [c for c in cues if c["dur"] > 0.7]


# def make_text_clip(text, width, font_path, fontsize, fill="#FFFFFF",
#                    stroke_fill="#000000", stroke_width=3, align="center", line_gap=8):
#     """Pillow text -> ImageClip with alpha (no ImageMagick)."""
#     try:
#         font = ImageFont.truetype(font_path, fontsize)
#     except Exception:
#         font = ImageFont.load_default()

#     margin = 16
#     W = int(width)
#     tmp = Image.new("RGBA", (W, 10), (0, 0, 0, 0))
#     tmp_draw = ImageDraw.Draw(tmp)

#     # wrap
#     words = text.split()
#     lines, line = [], ""
#     for w in words:
#         trial = (line + " " + w).strip()
#         if tmp_draw.textlength(trial, font=font) <= (W - 2*margin):
#             line = trial
#         else:
#             if line:
#                 lines.append(line)
#             line = w
#     if line:
#         lines.append(line)

#     bbox = font.getbbox("Ay")
#     line_h = (bbox[3] - bbox[1]) if bbox else fontsize + 6
#     H = (line_h + line_gap) * len(lines) + 2 * margin

#     img = Image.new("RGBA", (W, H), (0, 0, 0, 0))
#     draw = ImageDraw.Draw(img)

#     y = margin
#     for ln in lines:
#         w = draw.textlength(ln, font=font)
#         if align == "center":
#             x = int((W - w) // 2)
#         elif align == "left":
#             x = margin
#         else:
#             x = max(margin, W - margin - int(w))

#         if stroke_width > 0:
#             draw.text((x, y), ln, font=font, fill=stroke_fill,
#                       stroke_width=stroke_width, stroke_fill=stroke_fill)
#         draw.text((x, y), ln, font=font, fill=fill,
#                   stroke_width=stroke_width, stroke_fill=stroke_fill)
#         y += line_h + line_gap

#     arr = np.array(img)
#     if arr.shape[2] == 4:
#         rgb = arr[:, :, :3]
#         alpha = arr[:, :, 3] / 255.0
#         clip = ImageClip(rgb).set_mask(ImageClip(alpha, ismask=True))
#     else:
#         clip = ImageClip(arr)
#     return clip

def _autoshrink_fontsize_to_fit(text, width_px, font_path, start_size, min_size=36, margin=16):
    """Find largest fontsize that fits the text on ONE line within width_px."""
    try:
        size = start_size
        while size >= min_size:
            f = ImageFont.truetype(font_path, size)
            w = ImageDraw.Draw(Image.new("RGB", (10, 10))).textlength(text, font=f)
            if w <= (width_px - 2*margin):
                return size
            size -= 2
        return max(min_size, size)
    except Exception:
        return start_size

def make_text_clip_emoji_aware(text, width, font_path, fontsize, fill, stroke_fill, stroke_width, align):
    if EMOJI_RENDER_MODE == "skia":
        try:
            return make_text_clip_skia(text, width, font_path, EMOJI_FONT_PATH, fontsize, fill, stroke_fill, stroke_width, align)
        except Exception:
            pass  # fall back below
    # fallback to your existing PIL renderer (monochrome emoji)
    return make_text_clip(text, width, font_path, fontsize, fill, stroke_fill, stroke_width, align)

# If we can't render color emoji, remove VS16/ZWJ so no tofu boxes appear.
def _strip_emoji_variation_selectors(s: str) -> str:
    # VS15/VS16 + ZWJ cause tofu when the font lacks color emoji
    return (s or "").replace("\uFE0F", "").replace("\uFE0E", "").replace("\u200D", "")


def make_text_clip(text, width, font_path, fontsize, fill="#FFFFFF",
                   stroke_fill="#000000", stroke_width=3, align="center"):
    """
    Single-line super with auto-shrink to width.
    If skia + a color-emoji font are available, render with Skia so emoji appear in full color.
    Otherwise fall back to PIL (your previous behavior).
    """
    # --- Preprocess text & size (reuse your existing autoshrink) ---
    text = " ".join((text or "").strip().split())
    fontsize = _autoshrink_fontsize_to_fit(text, width, font_path, fontsize)

        # --- Try Skia path (full-color emoji) ---
    try:
        import skia
        from pathlib import Path

        emoji_font_path = EMOJI_FONT_PATH  # you already set this at the top
        use_skia = Path(emoji_font_path).exists()

        if use_skia:
            # Helpers
            def to_argb(hexcolor="#FFFFFF"):
                c = hexcolor.lstrip('#')
                if len(c) == 3:
                    c = "".join(ch * 2 for ch in c)
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return skia.Color4f(r/255, g/255, b/255, 1.0).toColor()

            # Load typefaces from files
            tf_main  = skia.Typeface.MakeFromFile(font_path)
            tf_emoji = skia.Typeface.MakeFromFile(emoji_font_path)

            # Register them in a provider/collection so Paragraph can resolve them
            provider = skia.TypefaceFontProvider()
            if tf_main:  provider.registerTypeface(tf_main)
            if tf_emoji: provider.registerTypeface(tf_emoji)

            collection = skia.FontCollection()
            collection.setAssetFontManager(provider)

            # Paragraph/Text style
            ps = skia.ParagraphStyle()
            ts = skia.TextStyle()
            fams = []
            if tf_main and tf_main.familyName():  fams.append(tf_main.familyName())
            if tf_emoji and tf_emoji.familyName(): fams.append(tf_emoji.familyName())
            if fams:
                ts.setFontFamilies(fams)
            ts.setFontSize(float(fontsize))
            ts.setColor(to_argb(fill))
            ps.setTextStyle(ts)
            ps.setTextAlign({
                "left":   skia.TextAlign.kLeft,
                "center": skia.TextAlign.kCenter,
                "right":  skia.TextAlign.kRight
            }[align])

            # Layout
            margin = 16
            max_w = float(int(width) - 2 * margin)
            builder = skia.ParagraphBuilder(ps, collection)
            builder.addText(text)
            p = builder.Build()
            p.layout(max_w)

            W = int(max_w + 2 * margin)
            H = int(max(32, p.height() + 2 * margin))

            # Surface
            surface = skia.Surface(W, H)
            canvas = surface.getCanvas()
            canvas.clear(skia.Color4f(0, 0, 0, 0).toColor())

            # Stroke (approximate): offset pass
            if stroke_width and stroke_width > 0:
                shadow_color = to_argb(stroke_fill)
                dx = dy = int(max(1, stroke_width))
                ps_shadow = skia.ParagraphStyle(ps)
                ts_shadow = skia.TextStyle(ts)
                ts_shadow.setColor(shadow_color)
                ps_shadow.setTextStyle(ts_shadow)

                for ox, oy in [(-dx, 0), (dx, 0), (0, -dy), (0, dy)]:
                    b2 = skia.ParagraphBuilder(ps_shadow, collection)
                    b2.addText(text)
                    p2 = b2.Build()
                    p2.layout(max_w)
                    p2.paint(canvas, ox, oy + margin)

            # Main fill pass
            p.paint(canvas, 0, margin)

            # Export → ImageClip with alpha
            img = surface.makeImageSnapshot()
            arr = np.frombuffer(img.tobytes(), dtype=np.uint8).reshape(H, W, 4)
            rgb = arr[:, :, :3]
            alpha = arr[:, :, 3] / 255.0
            return ImageClip(rgb).set_mask(ImageClip(alpha, ismask=True))
    except Exception:
        pass
    text = _strip_emoji_variation_selectors(text)

    # --- PIL fallback (previous behavior; emoji may be monochrome) ---
    try:
        font = ImageFont.truetype(font_path, fontsize)
    except Exception:
        font = ImageFont.load_default()

    margin = 16
    W = int(width)
    w = ImageDraw.Draw(Image.new("RGBA", (10, 10))).textlength(text, font=font)
    bbox = font.getbbox("Ay")
    line_h = (bbox[3] - bbox[1]) if bbox else fontsize + 6
    H = line_h + 2 * margin

    img = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    if align == "center":
        x = int((W - w) // 2)
    elif align == "left":
        x = margin
    else:
        x = max(margin, W - margin - int(w))
    y = margin

    if stroke_width > 0:
        draw.text((x, y), text, font=font, fill=stroke_fill,
                  stroke_width=stroke_width, stroke_fill=stroke_fill)
    draw.text((x, y), text, font=font, fill=fill,
              stroke_width=stroke_width, stroke_fill=stroke_fill)

    arr = np.array(img)
    if arr.shape[2] == 4:
        rgb = arr[:, :, :3]
        alpha = arr[:, :, 3] / 255.0
        clip = ImageClip(rgb).set_mask(ImageClip(alpha, ismask=True))
    else:
        clip = ImageClip(arr)
    return clip


def apply_text_rules(line: dict) -> dict:
    """
    Input: a cue dict (may include text, color).
    Output: possibly modified cue dict (text/color).
    Extend with client-specific rules as needed.
    """
    t = (line.get("text") or "").strip()

    # Cloud video: ensure a cloud emoji at the end of that specific sentence
    if "Come see its daily companions" in t and "☁️" not in t:
        t = t.rstrip() + " ☁️"

    # Special 'HI' color tweak (if they want a reminder style)
    if t == "HI" and not line.get("color"):
        line["color"] = "#FF61C0"  # adjust to brand

    line["text"] = t
    return line



# def render_video_clean(video_path: str, lines: List[str], cfg: dict) -> str:
#     """
#     Build one output video (clean style). Never drops text:
#     - readable cue planning
#     - 1080x1920 crop
#     - duration 20–25s target; extend via slow-mo/freezeframe if needed
#     - fade in/out
#     """
#     W, H = cfg["tiktok"]["width"], cfg["tiktok"]["height"]
#     margin = cfg["margin_px"]
#     font_path = cfg["font_path"]
#     font_size = cfg["font_size"]
#     text_color = cfg["text_color"]
#     stroke_color = cfg["stroke_color"]
#     stroke_width = cfg["stroke_width"]

#     target_min, target_max = 20.0, 25.0
#     fade_in, fade_out = 0.5, 1.0
#     EPS = 0.08
#     enforce_max = cfg.get("enforce_max", False)

#     base = VideoFileClip(video_path)

#     # fit vertical
#     target_ratio = W / H
#     src_ratio = base.w / base.h
#     if src_ratio > target_ratio:
#         new_w = int(base.h * target_ratio)
#         x1 = (base.w - new_w) // 2
#         base = base.crop(x1=x1, y1=0, x2=x1 + new_w, y2=base.h)
#     else:
#         new_h = int(base.w / target_ratio)
#         y1 = (base.h - new_h) // 2
#         base = base.crop(x1=0, y1=y1, x2=base.w, y2=y1 + new_h)
#     base = base.resize((W, H))

#     scenes = detect_scenes(video_path)
#     cues = plan_cues_readable(lines, scenes, base.duration, min_read=2.2, max_read=3.5, gap=0.12)

#     # compute needed length
#     last_end_needed = max((c["start"] + c["dur"] for c in cues), default=0.0) + 0.20
#     dur = base.duration
#     desired = max(target_min, last_end_needed, dur)
#     if enforce_max:
#         desired = min(desired, target_max)

#     # stretch via slow-mo
#     if dur < desired:
#         factor = dur / desired
#         if factor < 0.98:
#             base = base.fx(vfx.speedx, factor)
#             dur = base.duration

#     # pad with freeze frame if still short
#     if dur + 1e-3 < desired:
#         t_last = max(0.0, dur - 1.0 / (base.fps or 25.0))
#         frame = base.get_frame(t_last)
#         tail = ImageClip(frame).set_duration(desired - dur).crossfadein(0.25)
#         base = concatenate_videoclips([base, tail], method="compose", padding=-0.25)
#         dur = base.duration

#     # fades
#     base = base.fx(vfx.fadein, fade_in).fx(vfx.fadeout, fade_out)
#     if base.audio:
#         base = base.audio_fadein(fade_in).audio_fadeout(fade_out)

#     # overlay text
#     overlays = [base]
#     for cue in cues:
#         start = float(cue["start"])
#         d_c = float(cue["dur"])
#         if start >= dur - EPS:
#             start = max(0.0, dur - EPS - max(0.8, d_c))
#         end = min(dur - EPS, start + d_c)
#         d_c = max(0.8, end - start)

#         clip = make_text_clip(
#             cue["text"], width=int(W * 0.86),
#             font_path=font_path, fontsize=font_size,
#             fill=text_color, stroke_fill=stroke_color, stroke_width=stroke_width
#         )
#         pos_y = H - margin - clip.h if cue.get("pos", "bottom") == "bottom" else margin
#         overlays.append(
#             clip.set_start(start).set_duration(d_c).set_position(("center", pos_y))
#         )

#     out = CompositeVideoClip(overlays, size=(W, H)).set_duration(dur)
#     if base.audio:
#         out = out.set_audio(base.audio)

#     # output path
#     out_dir = Path(cfg["output_dir"])
#     out_dir.mkdir(parents=True, exist_ok=True)
#     stem = Path(video_path).stem
#     out_path = out_dir / f"{stem}_output.mp4"
#     out.write_videofile(
#         str(out_path),
#         fps=cfg["tiktok"]["fps"],
#         codec="libx264",
#         audio_codec="aac",
#         bitrate=cfg["tiktok"]["bitrate"],
#         verbose=False,
#         logger=None
#     )
#     return str(out_path)

def render_video_clean(video_path: str, cues_or_lines, cfg: dict) -> str:
    """
    Render one TikTok video.
    Accepts either:
      • a list of structured cues (dicts with optional start/dur/pos/font/color/size), or
      • a list of plain strings (legacy) → will auto-plan timing.
    """
    W, H = cfg["tiktok"]["width"], cfg["tiktok"]["height"]
    margin = cfg["margin_px"]
    base_font_path = cfg["font_path"]
    base_font_size = cfg["font_size"]
    default_color = cfg["text_color"]
    stroke_color = cfg["stroke_color"]
    stroke_width = cfg["stroke_width"]
    default_pos = cfg.get("default_pos", "middle")

    target_min, target_max = 20.0, 25.0
    fade_in, fade_out = 0.5, 1.0
    EPS = 0.08
    enforce_max = cfg.get("enforce_max", False)

    base = VideoFileClip(video_path)

    # fit vertical (9:16)
    target_ratio = W / H
    src_ratio = base.w / base.h
    if src_ratio > target_ratio:
        new_w = int(base.h * target_ratio)
        x1 = (base.w - new_w) // 2
        base = base.crop(x1=x1, y1=0, x2=x1 + new_w, y2=base.h)
    else:
        new_h = int(base.w / target_ratio)
        y1 = (base.h - new_h) // 2
        base = base.crop(x1=0, y1=y1, x2=base.w, y2=y1 + new_h)
    base = base.resize((W, H))

    # Build cues: either use explicit timing or auto-plan
    explicit = isinstance(cues_or_lines, list) and len(cues_or_lines) > 0 and isinstance(cues_or_lines[0], dict)
    # if explicit:
    #     cues = []
    #     for raw in cues_or_lines:
    #         c = dict(raw)  # copy
    #         c = apply_text_rules(c)  # apply client rules
    #         # compute dur if missing but end present
    #         if "dur" not in c and "start" in c and "end" in c:
    #             c["dur"] = max(0.05, float(c["end"]) - float(c["start"]))
    #         # fallback position
    #         c["pos"] = (c.get("pos") or default_pos or "middle").lower()
    #         cues.append(c)
    #     # sanity: discard negatives / fix sizes
    #     for c in cues:
    #         if "size" in c:
    #             try:
    #                 c["size"] = int(c["size"])
    #             except Exception:
    #                 c["size"] = base_font_size
    # else:
    #     # legacy path: list of strings -> auto time planner
    #     scenes = detect_scenes(video_path)
    #     cues = plan_cues_readable(cues_or_lines, scenes, base.duration, min_read=2.2, max_read=3.5, gap=0.12)
    #     # attach defaults
    #     for c in cues:
    #         c["pos"] = default_pos
    explicit = isinstance(cues_or_lines, list) and len(cues_or_lines) > 0 and isinstance(cues_or_lines[0], dict)
    if explicit:
        # Do any cues carry explicit timing?
        has_timing = any(("start" in c) or ("end" in c) or ("dur" in c) for c in cues_or_lines)

        if has_timing:
            cues = []
            for raw in cues_or_lines:
                c = dict(raw)
                c = apply_text_rules(c)

                # compute dur if end is provided
                if "dur" not in c and "start" in c and "end" in c:
                    c["dur"] = max(0.05, float(c["end"]) - float(c["start"]))

                # defaults for style
                c["pos"] = (c.get("pos") or default_pos or "middle").lower()
                if "size" in c:
                    try:
                        c["size"] = int(c["size"])
                    except Exception:
                        c["size"] = base_font_size
                # ensure color exists (keeps per-line color if given)
                c["color"] = c.get("color", default_color)

                cues.append(c)

        else:
            # 🔧 No timing in PPTX dicts → auto-plan timings but KEEP style (font/color/pos/size)
            # 1) get readable timings from scenes using only the text
            plain_lines = [apply_text_rules(dict(raw)).get("text", "").strip()
                        for raw in cues_or_lines if str(raw.get("text", "")).strip()]
            scenes = detect_scenes(video_path)
            planned = plan_cues_readable(plain_lines, scenes, base.duration, min_read=2.2, max_read=3.5, gap=0.12)

            # 2) merge timings back with original style fields, 1:1 in order
            cues = []
            for i, p in enumerate(planned):
                src = dict(cues_or_lines[i])  # corresponding original cue
                src = apply_text_rules(src)

                cue = {
                    "text": src.get("text", "").strip(),
                    "start": float(p["start"]),
                    "dur": float(p["dur"]),
                    # style with defaults
                    "pos": (src.get("pos") or default_pos or "middle").lower(),
                    "font": src.get("font"),                     # may be None → handled later
                    "size": int(src.get("size", base_font_size)),
                    "color": src.get("color", default_color),
                }
                cues.append(cue)
    else:
        # Legacy plain string mode
        scenes = detect_scenes(video_path)
        cues = plan_cues_readable(cues_or_lines, scenes, base.duration, min_read=2.2, max_read=3.5, gap=0.12)
        for c in cues:
            c["pos"] = default_pos
            c["color"] = default_color



    # Figure desired duration from cues (if explicit)
    dur = base.duration
    if cues:
        last_end_needed = 0.0
        for c in cues:
            s = float(c.get("start", 0.0))
            d = float(c.get("dur", 0.0)) if "dur" in c else 0.0
            last_end_needed = max(last_end_needed, s + d + 0.20)
        desired = max(target_min, last_end_needed, dur)
    else:
        desired = max(target_min, dur)

    if enforce_max:
        desired = min(desired, target_max)

    # stretch via slow-mo if short
    if dur < desired:
        factor = dur / desired
        if factor < 0.98:
            base = base.fx(vfx.speedx, factor)
            dur = base.duration

    # pad with freeze frame if still short
    if dur + 1e-3 < desired:
        t_last = max(0.0, dur - 1.0 / (base.fps or 25.0))
        frame = base.get_frame(t_last)
        tail = ImageClip(frame).set_duration(desired - dur).crossfadein(0.25)
        base = concatenate_videoclips([base, tail], method="compose", padding=-0.25)
        dur = base.duration

    # fades
    base = base.fx(vfx.fadein, fade_in).fx(vfx.fadeout, fade_out)
    if base.audio:
        base = base.audio_fadein(fade_in).audio_fadeout(fade_out)

    # overlay text
    overlays = [base]
    for cue in cues:
        # timings
        start = float(cue.get("start", 0.0))
        d_c = float(cue.get("dur", 0.0)) if "dur" in cue else 0.0
        if d_c <= 0.0:
            # if explicit time missing in explicit mode, show for a readable default
            d_c = 2.4
        if start >= dur - EPS:
            start = max(0.0, dur - EPS - max(0.8, d_c))
        end = min(dur - EPS, start + d_c)
        d_c = max(0.8, end - start)

        # style
        font_key = cue.get("font")  # 'default'|'cute'|'qq'|path
        this_font_path = resolve_font_path(font_key, cfg) if font_key else base_font_path
        this_size = int(cue.get("size", base_font_size))
        this_color = cue.get("color", default_color)

        # text (already rule-adjusted)
        text_final = cue.get("text", "").strip()
        if not text_final:
            continue

        clip = make_text_clip(
            text_final,
            width=int(W * 0.86),
            font_path=this_font_path,
            fontsize=this_size,
            fill=this_color,
            stroke_fill=stroke_color,
            stroke_width=stroke_width
        )

        # placement
        pos_key = (cue.get("pos") or default_pos).lower()
        if pos_key == "top":
            pos_y = margin
        elif pos_key == "middle":
            pos_y = int((H - clip.h) // 2)
        else:
            pos_y = H - margin - clip.h

        overlays.append(
            clip.set_start(start).set_duration(d_c).set_position(("center", pos_y))
        )

    out = CompositeVideoClip(overlays, size=(W, H)).set_duration(dur)
    if base.audio:
        out = out.set_audio(base.audio)

    # output path
    out_dir = Path(cfg["output_dir"])
    out_dir.mkdir(parents=True, exist_ok=True)
    stem = Path(video_path).stem
    out_path = out_dir / f"{stem}_output.mp4"
    out.write_videofile(
        str(out_path),
        fps=cfg["tiktok"]["fps"],
        codec="libx264",
        audio_codec="aac",
        bitrate=cfg["tiktok"]["bitrate"],
        verbose=False,
        logger=None
    )
    return str(out_path)


def save_uploaded_files(files, dest_dir: Path) -> List[str]:
    """Save uploaded video files to a folder and return paths."""
    paths = []
    for f in files:
        fp = dest_dir / f.name
        with open(fp, "wb") as out:
            out.write(f.read())
        paths.append(str(fp))
    return paths


def extract_zip_to_dir(zip_bytes: bytes, dest_dir: Path) -> List[str]:
    """Extract a zip to dest_dir, return list of video file paths discovered."""
    paths = []
    with zipfile.ZipFile(io.BytesIO(zip_bytes)) as z:
        z.extractall(dest_dir)
    for p in dest_dir.rglob("*"):
        if p.suffix.lower() in {".mp4", ".mov", ".m4v", ".avi", ".mkv"}:
            paths.append(str(p))
    return paths


# -----------------------------
# ----------- UI --------------
# -----------------------------
with st.sidebar:
    st.header("🎬 NP360 TikTok Generator")
    st.caption("Upload your **videos** (ZIP recommended) and the **instructions PPTX**.")

    vids_zip = st.file_uploader("Videos ZIP (preferred)", type=["zip"])
    vids_files = st.file_uploader("Or upload videos directly", type=["mp4", "mov", "m4v", "avi", "mkv"], accept_multiple_files=True)
    pptx_file = st.file_uploader("Instructions PPTX", type=["pptx"])

    st.markdown("**Guidelines**")
    st.markdown("""
- One slide **per video**  
- Slide **title = file name** (or unique substring)  
- Body: **line-by-line** supers (keep each ≤ 80 chars)  
- Final line starting with **`Copy:`** is the TikTok caption (not on-screen).  
- Keep style **simple & clean**, no colored frames.
    """, help="The agent maps videos to slides by finding the slide whose text contains the file name or substring.")

    st.markdown("---")
    st.markdown("**Rendering Options**")
    enforce = st.checkbox("Enforce 25s maximum video length", value=False,
                          help="If ON, videos are clamped to ≤ 25s (may compress last lines). If OFF, video will extend to fit all lines.")
    font = st.text_input("Font path (TTF/OTF)", value="assets/Fonts/Inter-VariableFont_opsz,wght.ttf")
    place = st.selectbox("Default super placement", ["top", "middle", "bottom"], index=1)
    cute_font = st.text_input("Cute font path (for 'font=cute')", value="assets/Fonts/Baloo2-Bold.ttf")
    qq_font = st.text_input("QQ style font path (for 'font=qq')", value="assets/Fonts/ComicNeue-Bold.ttf")

    size = st.slider("Font size", 40, 90, 70, 2)
    stroke = st.slider("Stroke width", 0, 6, 3, 1)

st.title("NP360 – TikTok Clean Version Generator")
st.markdown("<span class='badge'>POC</span> Generates one clean TikTok per input video, using the PPTX slide text.", unsafe_allow_html=True)
st.write("")

colA, colB = st.columns([1, 1])
with colA:
    st.subheader("1) Upload")
    st.markdown("Use the sidebar to upload **videos** and the **PPTX instructions**.")

with colB:
    st.subheader("2) Generate")
    run_btn = st.button("Generate Videos", type="primary", use_container_width=True)

st.markdown("---")

# -----------------------------
# ---------- RUN --------------
# -----------------------------
if run_btn:
    if (vids_zip is None and not vids_files) or (pptx_file is None):
        st.warning("Please upload videos (ZIP or files) **and** the instructions PPTX.")
    else:
        tmpdir = Path(tempfile.mkdtemp(prefix="np360_"))
        vids_dir = tmpdir / "videos"
        vids_dir.mkdir(parents=True, exist_ok=True)
        out_dir = tmpdir / "export"
        out_dir.mkdir(exist_ok=True)

        # Save PPTX
        pptx_path = tmpdir / "instructions.pptx"
        with open(pptx_path, "wb") as f:
            f.write(pptx_file.read())

        # Collect video paths
        video_paths = []
        if vids_zip is not None:
            video_paths.extend(extract_zip_to_dir(vids_zip.getvalue(), vids_dir))
        if vids_files:
            video_paths.extend(save_uploaded_files(vids_files, vids_dir))
        video_paths = sorted(video_paths)

        if not video_paths:
            st.error("No video files found. Please upload a zip of videos or add individual files.")
        else:
            cfg = {
                "tiktok": {"width": 1080, "height": 1920, "fps": 30, "bitrate": "10M"},
                "font_path": font,
                "font_size": size,
                "text_color": "#FFFFFF",
                "stroke_color": "#000000",
                "stroke_width": stroke,
                "margin_px": 120,
                "output_dir": str(out_dir),
                "enforce_max": enforce,
                "default_pos": place,
                "font_packs": {
                    "default": font,
                    "cute": cute_font,
                    "qq": qq_font
                },
            }

            st.info(f"Found **{len(video_paths)}** videos. Processing now…")
            progress = st.progress(0)
            status = st.empty()
            results = []

            for i, vp in enumerate(video_paths, start=1):
                fname = Path(vp).name
                status.write(f"Parsing slide and rendering: **{fname}**")
                # Map slide by searching the PPTX for the filename (stem)
                stem = Path(vp).stem
                # Try exact, then relaxed match (remove extension/underscores)
                # lines = extract_lines_from_pptx(str(pptx_path), stem)
                # if not lines:
                #     # Attempt relaxed match (e.g., only alnum)
                #     relaxed = re.sub(r"[^A-Za-z0-9]+", "", stem)
                #     lines = extract_lines_from_pptx(str(pptx_path), relaxed)

                # out_path = render_video_clean(vp, lines, cfg)
                res = extract_struct_from_pptx(str(pptx_path), stem, default_pos=cfg.get("default_pos", "middle"))
                if not res.cues:
                    # Attempt relaxed match (e.g., only alnum)
                    relaxed = re.sub(r"[^A-Za-z0-9]+", "", stem)
                    res = extract_struct_from_pptx(str(pptx_path), relaxed, default_pos=cfg.get("default_pos", "middle"))

                # Optionally use res.caption somewhere (e.g., print under preview)
                out_path = render_video_clean(vp, res.cues if res.cues else [], cfg)
                results.append(out_path)
                progress.progress(i / len(video_paths))

            status.write("✅ Finished generating all videos.")
            st.success(f"Generated **{len(results)}** videos.")

            st.markdown("---")
            st.subheader("Results")
            from pathlib import Path

            preview_cols = [1, 0.8, 1]  # middle column ~28% page width
            for p in results:
                with st.expander(Path(p).name, expanded=False):
                    _, mid, _ = st.columns(preview_cols)
                    with mid:
                        st.video(p) 
                        if res.caption:
                            st.markdown(f"<div class='small-muted'>Suggested caption: {res.caption}</div>", unsafe_allow_html=True)

                    with open(p, "rb") as f:
                        st.download_button("Download", f, file_name=Path(p).name, use_container_width=True)

            st.caption(f"Working directory: `{tmpdir}` (temporary). Move your downloads out if needed.")
